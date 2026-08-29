#!/usr/bin/env python3
"""Build the AQUA-SHIELD SIH 2026 Idea Submission deck -> PDF + PNG + PPTX.

FORMAT: follows the official SIH2026-IDEA-Presentation-Format template
(assets/SIH2026-IDEA-Presentation-Format.pptx) --
  * white ground, black serif ALL-CAPS slide titles
  * team-name oval top-left, SIH 2026 logo top-right
  * #0070C0 footer bar with "@SIH Idea submission- Template" and slide number
  * six slides including the title slide
  * the mandated content pointers on each slide are kept verbatim as headings,
    as the template's own instruction slide requires

Authored in HTML/CSS and rendered with a browser engine, which also produces
the PDF that the SIH portal accepts (the portal takes PDF only).

EVERY NUMBER IS EITHER:
  (a) MEASURED -- traceable to docs/BENCHMARKS.md, experiments/registry.jsonl,
      or a live run of the dashboard, or
  (b) TARGET   -- explicitly labelled as a goal, never as an achieved result.
Nothing is estimated or invented.

Usage:
    .venv/bin/python3 scripts/deck/build_deck.py
"""
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
DECK = ROOT / "scripts" / "deck"
IMG = ROOT / "docs" / "images"
OUT_DIR = ROOT / "docs"
BUILD = DECK / "_build"
BUILD.mkdir(exist_ok=True)

# ---------------------------------------------------------------------------
# EDIT THESE TWO LINES, then re-run this script.
TEAM_NAME = "Your Team Name"
TEAM_ID = "Your Team ID"
# ---------------------------------------------------------------------------

LOGO = (IMG / "sih_logo_2026.png").as_uri()
P = IMG / "prototype"


def u(name: str) -> str:
    return (P / name).as_uri()


def chrome(page_no: int, title: str, sub: str = "") -> str:
    subhtml = f'<span class="sub">{sub}</span>' if sub else ""
    return f"""
  <div class="sih-logo"><img src="{LOGO}" alt="Smart India Hackathon 2026"></div>
  <div class="team-oval">{TEAM_NAME}</div>
  <div class="slide-title">{title}{subhtml}</div>
  <div class="footbar">@SIH Idea submission- Template<span class="pg">{page_no}</span></div>"""


# ============================================================== SLIDE 1 =====
S1 = f"""
<div class="slide">
  <div class="sih-logo"><img src="{LOGO}" alt="Smart India Hackathon 2026"></div>

  <div style="position:absolute;left:60px;top:44px;">
    <div style="font-family:var(--serif);font-size:56px;font-weight:700;color:var(--navy);
                line-height:1.02;">SMART INDIA HACKATHON 2026</div>
  </div>

  <div style="position:absolute;left:60px;top:190px;width:830px;">
    <ul class="b" style="font-size:19.5px;">
      <li style="margin-bottom:26px;font-size:22px;">
        <b>Problem Statement ID</b> &ndash; <b>26057</b></li>
      <li style="margin-bottom:26px;">
        <b>Problem Statement Title</b> &ndash; AI-Powered Automated Underwater Marine
        Debris and Anomaly Detection System using Side-Scan Sonar Imagery</li>
      <li style="margin-bottom:26px;">
        <b>Theme</b> &ndash; Disaster Management
        <span style="color:var(--ink-3);">&nbsp;·&nbsp; Ministry of Earth Sciences (MoES) / NIOT</span></li>
      <li style="margin-bottom:26px;"><b>PS Category</b> &ndash; Software</li>
      <li style="margin-bottom:26px;"><b>Team ID</b> &ndash; {TEAM_ID}</li>
      <li><b>Team Name (Registered on portal)</b> &ndash; {TEAM_NAME}</li>
    </ul>
  </div>

  <div style="position:absolute;left:60px;bottom:44px;width:880px;">
    <div class="box teal t-teal" style="padding:13px 17px;">
      <span style="font-size:18px;font-weight:700;color:var(--teal);">AQUA-SHIELD</span>
      <span style="font-size:15.5px;color:var(--ink-2);"> &nbsp;&mdash;&nbsp; Detect &rarr; Verify
      &rarr; Localize &rarr; Act.&nbsp; A working, offline side-scan sonar pipeline that
      does not just detect, but <b>verifies every detection against physical acoustic
      evidence</b> before an operator ever sees it.</span>
    </div>
  </div>

  <div style="position:absolute;right:54px;top:215px;width:565px;">
    <img src="{u('panel_sonar.png')}"
         style="width:100%;border:1.6px solid var(--blue);border-radius:6px;display:block;">
    <div style="font-size:13px;color:var(--ink-3);margin-top:8px;line-height:1.35;
                text-align:center;">
      Live output of our running prototype on a real held-out survey frame
      (<span class="mono">0460_2018</span>).<br>Yellow = man-made, green = ambiguous.
    </div>
  </div>
</div>"""


# ============================================================== SLIDE 2 =====
S2 = f"""
<div class="slide">
  {chrome(2, "AQUA-SHIELD",
          "Acoustic Intelligence for Underwater Anomaly, Debris &amp; Marine-Hazard Localization")}
  <div class="body">
    <div class="fill" style="display:grid;grid-template-columns:1.08fr 1fr;gap:14px;">

      <div class="box" style="display:flex;flex-direction:column;">
        <h2 class="ptr">Proposed Solution</h2>
        <ul class="b">
          <li><b>Ten-stage offline pipeline.</b> Ingest &rarr; quality control &rarr;
              tiling &rarr; <span class="hb">YOLO11n detection</span> &rarr;
              <span class="hi">physical verification</span> &rarr; calibration &rarr;
              de-duplication &rarr; geolocation &rarr; priority &rarr; report.</li>
          <li><b>Independent verification stage.</b> Every candidate is re-examined
              against <span class="hi">10 features measured from the pixels</span> &mdash;
              shadow coherence, local SNR, contrast, compactness, texture &mdash; that do
              <i>not</i> depend on the detector's own opinion.</li>
          <li><b>Calibrated confidence.</b> Platt scaling fitted on a held-out survey;
              when unfitted, every hazard is stamped
              <span class="mono ha">calibrated: false</span>.</li>
          <li><b>Ranked hazard register.</b> The output is not a model file &mdash; it is a
              prioritised, auditable list exported as
              <span class="mono">GeoJSON / CSV / SQLite</span>, opening directly in QGIS.</li>
          <li><b>Runs fully offline on a laptop.</b> No cloud, no per-image cost, no data
              leaving the vessel.</li>
        </ul>
        <div style="margin-top:12px;">
          <div style="font-size:12.5px;font-weight:700;color:var(--ink-2);margin-bottom:6px;">
            THE TEN PHYSICAL FEATURES THE VERIFIER MEASURES</div>
          <span class="tag">shadow ratio</span><span class="tag">shadow side consistency</span>
          <span class="tag">local SNR</span><span class="tag">contrast</span>
          <span class="tag">compactness</span><span class="tag">aspect ratio</span>
          <span class="tag">edge density</span><span class="tag">relative texture</span>
          <span class="tag">area fraction</span><span class="tag">range position</span>
        </div>

        <div style="margin-top:14px;">
          <div style="font-size:12.5px;font-weight:700;color:var(--ink-2);margin-bottom:7px;">
            THE TEN STAGES &mdash; EACH SEPARATELY TESTABLE AND SEPARATELY ABLATED</div>
          <div class="row"><div class="step " style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Ingest</div></div><div class="arw" style="flex:0 0 8px;font-size:11px;">&rsaquo;</div><div class="step " style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Quality<br>control</div></div><div class="arw" style="flex:0 0 8px;font-size:11px;">&rsaquo;</div><div class="step off" style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Pre-<br>process</div></div><div class="arw" style="flex:0 0 8px;font-size:11px;">&rsaquo;</div><div class="step " style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Tiling</div></div><div class="arw" style="flex:0 0 8px;font-size:11px;">&rsaquo;</div><div class="step " style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Detection</div></div><div class="arw" style="flex:0 0 8px;font-size:11px;">&rsaquo;</div><div class="step key" style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Verify</div></div><div class="arw" style="flex:0 0 8px;font-size:11px;">&rsaquo;</div><div class="step " style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Calibrate</div></div><div class="arw" style="flex:0 0 8px;font-size:11px;">&rsaquo;</div><div class="step " style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Dedup</div></div><div class="arw" style="flex:0 0 8px;font-size:11px;">&rsaquo;</div><div class="step " style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Geo-<br>locate</div></div><div class="arw" style="flex:0 0 8px;font-size:11px;">&rsaquo;</div><div class="step " style="padding:5px 3px;"><div class="t" style="font-size:10px;line-height:1.15;">Report</div></div></div>
          <div style="font-size:11px;color:var(--ink-3);margin-top:7px;line-height:1.35;">
            <span style="color:var(--teal);font-weight:700;">Teal</span> = our core
            contribution. &nbsp;<span style="color:var(--amber);font-weight:700;">Dashed
            amber</span> = built, measured, and <b>switched off by default because it made
            detection worse</b> &mdash; we keep negative results rather than hiding them.
          </div>
        </div>

        <div class="note b" style="margin-top:auto;">
          <b>Status: a working prototype, not a concept.</b> It is trained, measured on
          held-out surveys, covered by 117 passing tests &mdash; and every screenshot in
          this deck is its live output.
        </div>
      </div>

      <div style="display:flex;flex-direction:column;gap:12px;">
        <div class="box red" style="flex:1;display:flex;flex-direction:column;">
          <h2 class="ptr red">How It Addresses the Problem</h2>
          <div style="display:flex;align-items:center;gap:15px;margin-bottom:9px;">
            <div style="font-size:54px;font-weight:800;color:var(--red);line-height:1;">74%</div>
            <div style="font-size:14px;line-height:1.35;color:var(--ink-2);">
              of frames in our held-out surveys contain <b>no target at all</b>
              (473 of 612 &mdash; measured). The real difficulty is not finding objects, it
              is <b>not crying wolf on empty seabed</b>.</div>
          </div>
          <ul class="b tight">
            <li>A recall-tuned detector alarms constantly, the analyst stops trusting it,
                and the system is abandoned. <span class="hi">Precision is the binding
                constraint</span> &mdash; so we built a stage specifically to enforce it.</li>
            <li>Replaces manual reading of thousands of km of sonar with a ranked register
                the operator can audit, accept or reject.</li>
            <li>Seabed clutter &mdash; rocks, sand ripples, acoustic shadows &mdash; is
                rejected on physical evidence, not on a hand-tuned threshold.</li>
          </ul>
          <div style="margin-top:auto;padding-top:10px;display:flex;align-items:center;gap:15px;">
            <div style="display:grid;grid-template-columns:repeat(20,11px);gap:3px;">
              <span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#D8DFE6;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span><span style="width:11px;height:11px;border-radius:2px;display:inline-block;background:#C00000;"></span>
            </div>
            <div style="font-size:11.5px;color:var(--ink-3);line-height:1.35;">
              <span style="color:#C00000;font-weight:700;">&#9632;</span> frames with a target<br>
              <span style="color:#B6BEC6;font-weight:700;">&#9632;</span> empty seabed<br>
              <span style="font-size:10.5px;">100 frames, to scale</span>
            </div>
          </div>
        </div>

        <div class="box teal" style="flex:1;">
          <h2 class="ptr teal">Innovation &amp; Uniqueness</h2>
          <ul class="b tight">
            <li><b>Verification on independent physical evidence</b> &mdash; not a second
                neural network agreeing with the first.</li>
            <li><b>Refusal as a feature.</b> No navigation metadata &rarr; no coordinate.
                We return <span class="mono ha">null</span> plus the reason. A fabricated
                latitude exports cleanly to CSV and sends a vessel to open water.</li>
            <li><b>Confidence &ne; Priority.</b> "Is it real?" and "should you care?" are
                scored separately, so a low-confidence large hazard is not buried.</li>
            <li><b>The verifier doubles as a diagnostic instrument.</b> Its fitted weights
                gave acoustic shadow a <i>negative</i> weight &mdash; the opposite of the
                physics. That exposed a real train/inference preprocessing mismatch in our
                own pipeline; fixing it moved F1
                <span class="mono hg">0.012 &rarr; 0.144</span>. A hand-tuned threshold
                would have hidden the defect.</li>
          </ul>
        </div>
      </div>
    </div>
  </div>
</div>"""


# ============================================================== SLIDE 3 =====
def step(t, d, cls=""):
    return f'<div class="step {cls}"><div class="t">{t}</div><div class="d">{d}</div></div>'


A = '<div class="arw">&rsaquo;</div>'

S3 = f"""
<div class="slide">
  {chrome(3, "TECHNICAL APPROACH")}
  <div class="body">
    <div class="fill" style="display:grid;grid-template-columns:0.60fr 1fr;gap:14px;">

      <div style="display:flex;flex-direction:column;gap:11px;">
        <div class="box violet" style="flex:1;display:flex;flex-direction:column;">
          <h3 class="sec" style="color:var(--violet);">Technologies to be Used</h3>
          <div style="font-size:12.5px;font-weight:700;color:var(--ink-2);margin:6px 0 4px;">
            Language &amp; machine learning</div>
          <div><span class="tag v">Python 3.12</span><span class="tag v">PyTorch 2.13</span>
               <span class="tag v">Ultralytics YOLO11n</span><span class="tag v">scikit-learn</span>
               <span class="tag v">ONNX Runtime</span></div>
          <div style="font-size:12.5px;font-weight:700;color:var(--ink-2);margin:9px 0 4px;">
            Sonar &amp; geospatial</div>
          <div><span class="tag">OpenCV</span><span class="tag">NumPy</span>
               <span class="tag">pyproj</span><span class="tag">GeoTIFF</span>
               <span class="tag">GeoJSON &rarr; QGIS</span></div>
          <div style="font-size:12.5px;font-weight:700;color:var(--ink-2);margin:9px 0 4px;">
            Application &amp; quality</div>
          <div><span class="tag g">Streamlit console</span><span class="tag g">FastAPI REST</span>
               <span class="tag g">SQLite register</span><span class="tag g">pytest &middot; 117 tests</span></div>
          <div style="font-size:12.5px;font-weight:700;color:var(--ink-2);margin:9px 0 4px;">
            Compute targets</div>
          <div><span class="tag a">Apple MPS</span><span class="tag a">CUDA</span>
               <span class="tag a">CPU fallback</span><span class="tag a">Jetson-class edge (target)</span></div>

          <div style="margin-top:auto;padding-top:13px;">
            <div style="font-size:12.5px;font-weight:700;color:var(--ink-2);margin-bottom:6px;">
              KEY TECHNICAL DECISIONS &mdash; AND WHY</div>
            <table class="t" style="font-size:11.5px;">
              <tr><td style="width:38%;"><b>YOLO11n</b>, not a<br>transformer detector</td>
                  <td class="dim">2.58 M params fits edge hardware; reviewed DETR variants
                      cost 2.6&times;&ndash;44&times; more for single-digit AP gains.</td></tr>
              <tr><td><b>Logistic verifier</b>,<br>not a second CNN</td>
                  <td class="dim">Inspectable weights, per-detection attribution &mdash; an
                      operator can see <i>why</i> a candidate was rejected.</td></tr>
            </table>
          </div>
        </div>

        <div class="box green">
          <h3 class="sec green">Data &mdash; real, licensed, leakage-free</h3>
          <table class="t" style="font-size:12.5px;">
            <tr><th>Dataset</th><th>Scale</th><th>Split by</th></tr>
            <tr><td>MILCO / NOMBO<br><span style="font-size:11px;color:var(--green);">CC BY 4.0</span></td>
                <td class="mono">465 / 93 / 612 frames<br>
                    <span style="font-size:10.5px;color:var(--ink-3);">191 test objects</span></td>
                <td class="dim">acquisition<br>year</td></tr>
            <tr><td>Derelict crab pot<br><span style="font-size:11px;color:var(--green);">CC BY-SA 4.0</span></td>
                <td class="mono">6,674 images<br>
                    <span style="font-size:10.5px;color:var(--ink-3);">9,311 objects &middot; 107 recordings</span></td>
                <td class="dim">recording<br>ID</td></tr>
          </table>
          <div class="note" style="margin-top:9px;padding:7px 11px;font-size:12px;">
            <b>Never a random split.</b> Consecutive sonar frames share seabed, gain
            settings and often the same object &mdash; a random split leaks and inflates
            every number. A unit test asserts the splits stay disjoint.
          </div>
        </div>
      </div>

      <div class="box" style="display:flex;flex-direction:column;">
        <h3 class="sec">Methodology &amp; Process for Implementation</h3>

        <div class="lane" style="margin-top:7px;">
          <div class="name">1 &nbsp;&middot;&nbsp; Acquisition &amp; conditioning</div>
          <div class="row">
            {step("Ingest", "SSS frames + nav CSV /<br>GeoTIFF")}{A}
            {step("Quality control", "dynamic range, speckle index,<br>dropout, water column")}{A}
            {step("Preprocess", "OFF by default &mdash; measured<br>3 ways, it hurt every time", "off")}{A}
            {step("Tiling", "640 px overlap, seam merge<br>by IoU / IoS")}
          </div>
        </div>

        <div class="lane teal" style="margin-top:17px;">
          <div class="name">2 &nbsp;&middot;&nbsp; Detection &amp; verification &mdash; the core contribution</div>
          <div class="row">
            {step("Detection", "YOLO11n &middot; 2.58 M params<br>6.3 GFLOPs")}{A}
            {step("10 physical features", "shadow ratio &middot; shadow side &middot; SNR &middot;<br>contrast &middot; compactness &middot; texture &hellip;", "key")}{A}
            {step("Logistic verifier", "L2-regularised, fitted on a<br>held-out survey", "key")}{A}
            {step("Calibration", "Platt &rarr; probability, or<br>stamped uncalibrated")}
          </div>
        </div>

        <div class="lane green" style="margin-top:17px;">
          <div class="name">3 &nbsp;&middot;&nbsp; Fusion, localization &amp; action</div>
          <div class="row">
            {step("De-duplication", "repeat sightings &rarr;<br>one unique hazard")}{A}
            {step("Geolocation", "GeoTIFF affine / per-ping nav<br><b>or refuse</b>")}{A}
            {step("Priority scoring", "separate from confidence")}{A}
            {step("Report", "ranked register &rarr;<br>GeoJSON &middot; CSV &middot; SQLite")}
          </div>
        </div>

        <div class="box t-grey grey" style="margin-top:17px;padding:9px 13px;">
          <div style="display:flex;align-items:center;gap:8px;flex-wrap:wrap;font-size:12.5px;">
            <b style="color:var(--navy);">I/O CONTRACT</b>
            <span class="mono">SSS frame 1024&times;1024 + nav</span><span style="color:var(--blue);">&rsaquo;</span>
            <span class="mono">640 px tiles</span><span style="color:var(--blue);">&rsaquo;</span>
            <span class="mono">boxes + class + raw score</span><span style="color:var(--blue);">&rsaquo;</span>
            <span class="mono" style="color:var(--teal);font-weight:700;">10-feature vector &rarr; p(real)</span>
            <span style="color:var(--blue);">&rsaquo;</span>
            <span class="mono">hazard record</span><span style="color:var(--blue);">&rsaquo;</span>
            <span class="mono">GeoJSON &middot; CSV</span>
          </div>
          <div style="font-size:11.5px;color:var(--ink-3);margin-top:6px;line-height:1.35;">
            <b>hazard record</b> = <span class="mono">{{id, class, level1, confidence,
            calibrated, priority, priority_band, lat, lon, uncertainty_m, observations,
            evidence[10], provenance}}</span> &mdash; every field an operator needs to act
            on, audit, or reject the call.
          </div>
        </div>

        <div style="margin-top:15px;flex:1;min-height:0;display:flex;flex-direction:column;">
          <h3 class="sec teal" style="margin-bottom:6px;">Working prototype &mdash; the operator console</h3>
          <div style="display:grid;grid-template-columns:1fr 1.42fr;gap:11px;flex:1;min-height:0;">
            <div style="display:flex;flex-direction:column;min-height:0;">
              <img src="{u('panel_stats.png')}"
                   style="display:block;width:100%;border:1px solid var(--rule);border-radius:4px;">
              <div class="cap">Run summary: 8 frames &rarr; 2 raw candidates &rarr;
                2 unique hazards &rarr; 1 HIGH priority.</div>
            </div>
            <div style="display:flex;flex-direction:column;min-height:0;">
              <img src="{u('panel_register.png')}"
                   style="display:block;width:100%;border:1px solid var(--rule);border-radius:4px;">
              <div class="cap">The hazard register. This survey ships no navigation, so
                <span class="mono">lat / lon / &plusmn;m</span> come back
                <span class="mono">None</span> &mdash; the system
                <b>refuses to invent a coordinate</b> rather than guessing.</div>
            </div>
          </div>
        </div>
      </div>
    </div>
  </div>
</div>"""


# ============================================================== SLIDE 4 =====
S4 = f"""
<div class="slide">
  {chrome(4, "FEASIBILITY AND VIABILITY")}
  <div class="body">
    <div class="fill" style="display:grid;grid-template-columns:1fr 1.04fr;gap:14px;">

      <div style="display:flex;flex-direction:column;gap:11px;">
        <div class="box teal" style="flex:1;display:flex;flex-direction:column;">
          <h2 class="ptr teal">Analysis of the Feasibility</h2>
          <div style="font-size:13.5px;color:var(--ink-2);margin-bottom:9px;line-height:1.35;">
            <b>It already runs.</b> A live run on 8 held-out frames:
            8 frames &rsaquo; 2 raw candidates &rsaquo; verify &rsaquo; 2 unique hazards
            &rsaquo; 1 HIGH priority, at 27 ms per frame.
          </div>
          <div style="display:grid;grid-template-columns:0.9fr 1fr;gap:10px;flex:1;min-height:0;">
            <div>
              <img src="{u('panel_sonar.png')}"
                   style="display:block;width:100%;border:1px solid var(--rule);border-radius:4px;">
              <div class="cap">Live prototype output, frame
                <span class="mono">0460_2018</span> &mdash; not a mockup.</div>
            </div>
            <div>
              <table class="t" style="font-size:12.5px;">
                <tr><th colspan="2">Measured on this machine</th></tr>
                <tr><td>MPS inference</td><td class="meas">21.4 ms</td></tr>
                <tr><td>ONNX export</td><td class="meas">10.6 MB &middot; 8.5 ms</td></tr>
                <tr><td>Throughput</td><td class="meas">37.4 frames/s</td></tr>
                <tr><td>Peak memory</td><td class="meas">640 MB</td></tr>
                <tr><td>Test suite</td><td class="meas">117 passing</td></tr>
              </table>
              <div class="cap" style="margin-top:6px;">Apple M5, 24 GB &mdash; a laptop, not
                a cluster. Reproducible from <span class="mono">experiments/</span>.</div>
            </div>
          </div>
        </div>

        <div class="box">
          <h3 class="sec">Verification stage &mdash; measured effect (612 held-out frames, 473 empty)</h3>
          <table class="t" style="font-size:12.5px;">
            <tr><th>Configuration</th><th class="c">Precision</th><th class="c">True positives</th>
                <th class="c">Falsely-alarmed frames</th></tr>
            <tr><td>Detector only</td><td class="c mono">0.247</td><td class="c mono">21</td>
                <td class="c mono">37 / 473</td></tr>
            <tr><td>+ hand-written rules</td><td class="c mono">0.300</td><td class="c mono">12</td>
                <td class="c mono">18 / 473</td></tr>
            <tr class="ours"><td><b>+ learned verification</b></td><td class="c meas">0.322</td>
                <td class="c meas">19</td><td class="c meas">25 / 473</td></tr>
          </table>
          <div class="cap" style="margin-top:5px;">Hand-written rules buy quiet by
            discarding targets (21 &rarr; 12). The learned verifier keeps <b>19 of 21</b>
            while cutting falsely-alarmed frames by <b>32%</b>. That gap is the case for
            fitting the filter rather than tuning thresholds.</div>
        </div>
      </div>

      <div style="display:flex;flex-direction:column;gap:11px;">
        <div class="box green" style="flex:1;">
          <h3 class="sec green">Target validation &mdash; what we will prove at the finals</h3>
          <table class="t" style="font-size:12.5px;">
            <tr><th style="width:42%;">Metric that matters</th>
                <th style="width:31%;">Measured today</th><th style="width:27%;">Target</th></tr>
            <tr><td>False alarms on target-free frames</td>
                <td class="meas">25 / 473 (5.3%)</td><td class="tgt">&le; 2%</td></tr>
            <tr><td>Precision @ IoU 0.3</td><td class="meas">0.322</td><td class="tgt">&ge; 0.60</td></tr>
            <tr><td>Recall @ IoU 0.3</td><td class="meas">0.099</td><td class="tgt">&ge; 0.45</td></tr>
            <tr><td>Ghost-gear detection, mAP50</td>
                <td class="meas">0.323 <span style="color:var(--ink-3);font-size:10.5px;">raw, 1 run</span></td>
                <td class="tgt">&ge; 0.55</td></tr>
            <tr><td>Calibration error (ECE)</td><td class="dim">fitted, not yet reported</td>
                <td class="tgt">&le; 0.05</td></tr>
            <tr><td>Geolocation error (CEP)</td>
                <td class="dim">not measurable &mdash; no nav<br>in available data</td>
                <td class="tgt">&le; 10 m</td></tr>
            <tr><td>Inference latency / frame</td><td class="meas">21.4 ms MPS</td>
                <td class="tgt">&le; 30 ms edge</td></tr>
          </table>
          <div class="note g" style="margin-top:9px;padding:7px 11px;font-size:12px;">
            <b>Measured</b> = reproducible from <span class="mono">experiments/</span> on
            held-out surveys. <b>Target</b> = a goal, not a result. We publish both columns
            so the gap is visible rather than hidden.
            <b>Baselines we will measure against:</b> detector-only &middot; hand-written
            rule filter &middot; YOLO11s / RT-DETR at equal compute.
          </div>
        </div>

        <div class="box amber">
          <h3 class="sec amber">Potential challenges &amp; risks &rarr; strategies for overcoming them</h3>
          <table class="t" style="font-size:12.5px;">
            <tr><th style="width:38%;">Challenge / risk</th><th>Strategy &mdash; already built</th></tr>
            <tr><td>Scarce labelled coastal SSS data</td>
                <td class="dim">Leakage-free recording-level splits; two licensed datasets;
                    honest reporting instead of inflated splits.</td></tr>
            <tr><td>Detection degrades under speckle noise</td>
                <td class="dim">Noise-robust checkpoint trained and shipped separately,
                    labelled as a trade-off rather than silently swapped in.</td></tr>
            <tr><td>Detector backend is AGPL-3.0</td>
                <td class="dim">Isolated behind an interface; the ONNX export path is
                    backend-independent.</td></tr>
            <tr><td>Edge compute is constrained on an AUV</td>
                <td class="dim">10.6 MB ONNX model measured at 8.5 ms &mdash; sized for
                    Jetson-class hardware.</td></tr>
            <tr><td>Missing navigation metadata</td>
                <td class="dim">The system refuses to emit a coordinate and states why,
                    rather than guessing.</td></tr>
          </table>
        </div>
      </div>
    </div>
  </div>
</div>"""


# ============================================================== SLIDE 5 =====
S5 = f"""
<div class="slide">
  {chrome(5, "IMPACT AND BENEFITS")}
  <div class="body">
    <div class="fill" style="display:grid;grid-template-columns:1fr 1fr;gap:14px;">

      <div style="display:flex;flex-direction:column;gap:11px;">
        <div class="box">
          <h2 class="ptr">Potential Impact on the Target Audience</h2>
          <div style="display:flex;align-items:stretch;gap:8px;margin:2px 0 11px;">
            <div class="box red t-red" style="flex:1;text-align:center;padding:9px 7px;">
              <div style="font-size:11.5px;font-weight:700;color:var(--red);">TODAY</div>
              <div style="font-size:12.5px;margin-top:3px;line-height:1.25;">analyst scrolls<br>raw sonar imagery</div>
            </div>
            <div class="arw" style="align-self:center;font-size:19px;">&rsaquo;</div>
            <div class="box teal t-teal" style="flex:1.3;text-align:center;padding:9px 7px;">
              <div style="font-size:11.5px;font-weight:700;color:var(--teal);">WITH AQUA-SHIELD</div>
              <div style="font-size:12.5px;margin-top:3px;line-height:1.25;">ranked hazard register,<br>evidence attached</div>
            </div>
            <div class="arw" style="align-self:center;font-size:19px;">&rsaquo;</div>
            <div class="box green t-green" style="flex:1;text-align:center;padding:9px 7px;">
              <div style="font-size:11.5px;font-weight:700;color:var(--green);">ACTION</div>
              <div style="font-size:12.5px;margin-top:3px;line-height:1.25;">GeoJSON &rarr; QGIS,<br>vessel tasked</div>
            </div>
          </div>
          <table class="t" style="font-size:12.5px;">
            <tr><th style="width:33%;">Beneficiary</th><th>Impact</th></tr>
            <tr><td><b>MoES / NIOT</b><br><span style="font-size:10.5px;color:var(--ink-3);">survey &amp; cleanup operations</span></td>
                <td class="dim">Direct fit to the problem statement's own agency; runs
                    offline on survey hardware already aboard.</td></tr>
            <tr><td><b>Fisheries &amp; coastal management</b></td>
                <td class="dim">Ghost gear keeps killing catch and habitat long after it is
                    lost; early detection breaks that cycle.</td></tr>
            <tr><td><b>Ports &amp; navigation safety</b></td>
                <td class="dim">Mine-like and man-made bottom objects flagged and
                    prioritised before they foul a channel.</td></tr>
            <tr><td><b>Environmental researchers</b></td>
                <td class="dim">Auditable, licence-clean datasets and a reproducible
                    experiment registry.</td></tr>
          </table>
        </div>

        <div class="box violet">
          <h3 class="sec" style="color:var(--violet);">Deployment &amp; scalability &mdash; offline-first by design</h3>
          <ul class="b tight">
            <li><b>Today:</b> runs on a single laptop (Apple M5, 24 GB), fully offline.</li>
            <li><b>Edge:</b> 10.6 MB ONNX export measured at 8.5 ms &mdash; sized for a
                Jetson aboard an AUV or survey launch.</li>
            <li><b>Fleet:</b> FastAPI + SQLite register scales the same pipeline from one
                operator console to a shore-side survey fleet.</li>
            <li><span class="ha"><b>Not yet:</b></span> never run on a Jetson or a live AUV.
                That is the next milestone, and we state it as a target rather than an
                achievement.</li>
          </ul>
        </div>
      </div>

      <div style="display:flex;flex-direction:column;gap:11px;">
        <div class="box green">
          <h2 class="ptr green">Benefits of the Solution</h2>
          <table class="t" style="font-size:12.5px;">
            <tr><th style="width:23%;">Dimension</th><th>Benefit</th></tr>
            <tr><td><b style="color:var(--green);">Environmental</b></td>
                <td class="dim">Faster location of ghost nets and derelict gear that keep
                    entangling marine life and smothering reefs after they are lost.</td></tr>
            <tr><td><b style="color:var(--blue);">Economic</b></td>
                <td class="dim">No cloud inference, no per-image API fee, no data egress
                    &mdash; the marginal cost of one more survey is the electricity to run
                    a laptop already aboard.</td></tr>
            <tr><td><b style="color:var(--violet);">Operational</b></td>
                <td class="dim">32% fewer falsely-alarmed frames while keeping 19 of 21 true
                    positives (measured) &mdash; the analyst reviews a ranked register
                    instead of raw imagery.</td></tr>
            <tr><td><b style="color:var(--teal);">Safety</b></td>
                <td class="dim">Mine-like and man-made hazards prioritised before a vessel
                    or diver encounters them.</td></tr>
            <tr><td><b style="color:var(--amber);">Social &amp;<br>institutional</b></td>
                <td class="dim">Every hazard carries evidence, calibrated confidence and
                    provenance, so a decision can be audited &mdash; not merely trusted.</td></tr>
          </table>
        </div>

        <div class="box" style="flex:1;display:flex;flex-direction:column;">
          <h3 class="sec">The exported hazard register &mdash; from the running system</h3>
          <img class="shot" src="{u('panel_geo_table.png')}">
          <div class="cap">Real coordinates, per-hazard uncertainty in metres, and a
            priority band &mdash; the artefact a cleanup vessel is tasked from.
            <b>The navigation track in this demo is synthetic and the product says so on
            screen:</b> the geolocation path works end to end; the fix is not validated.</div>
          <div class="note" style="margin-top:9px;padding:7px 12px;font-size:12.5px;">
            <b>What we deliberately do not claim.</b> No "% of analyst time saved" appears
            on this slide &mdash; that needs a user study we have not run.
          </div>
          <div style="margin-top:auto;padding-top:11px;display:grid;
                      grid-template-columns:repeat(3,1fr);gap:9px;">
            <div class="box teal t-teal" style="padding:9px 11px;text-align:center;">
              <div style="font-size:24px;font-weight:800;color:var(--teal);line-height:1;">&minus;32%</div>
              <div style="font-size:11px;color:var(--ink-2);margin-top:5px;line-height:1.28;">
                falsely-alarmed frames<br><span style="color:var(--ink-3);">measured</span></div>
            </div>
            <div class="box teal t-teal" style="padding:9px 11px;text-align:center;">
              <div style="font-size:24px;font-weight:800;color:var(--teal);line-height:1;">19/21</div>
              <div style="font-size:11px;color:var(--ink-2);margin-top:5px;line-height:1.28;">
                true positives kept<br><span style="color:var(--ink-3);">measured</span></div>
            </div>
            <div class="box teal t-teal" style="padding:9px 11px;text-align:center;">
              <div style="font-size:24px;font-weight:800;color:var(--teal);line-height:1;">37.4/s</div>
              <div style="font-size:11px;color:var(--ink-2);margin-top:5px;line-height:1.28;">
                frames processed<br><span style="color:var(--ink-3);">measured</span></div>
            </div>
          </div>
        </div>
      </div>
    </div>
  </div>
</div>"""


# ============================================================== SLIDE 6 =====
S6 = f"""
<div class="slide">
  {chrome(6, "RESEARCH AND REFERENCES")}
  <div class="body">
    <div class="fill" style="display:grid;grid-template-columns:1fr 1.02fr;gap:14px;">

      <div style="display:flex;flex-direction:column;gap:11px;">
        <div class="box">
          <h3 class="sec">Datasets &amp; prior art &mdash; every licence verified</h3>
          <table class="t" style="font-size:12px;">
            <tr><th style="width:35%;">Work</th><th style="width:20%;">Licence</th><th>Relationship</th></tr>
            <tr><td><b>MILCO / NOMBO</b><br><span style="font-size:10px;color:var(--ink-3);">Data in Brief 53:110132 &middot;<br>DOI 10.6084/m9.figshare.24574879</span></td>
                <td style="color:var(--green);font-weight:700;">CC BY 4.0</td>
                <td class="dim"><b>Our training data</b> &mdash; mine-like &amp; bottom objects.</td></tr>
            <tr><td><b>sss-crab-pot-detection-ds</b><br><span style="font-size:10px;color:var(--ink-3);">PINGEcosystem &middot; HuggingFace</span></td>
                <td style="color:var(--green);font-weight:700;">CC BY-SA 4.0</td>
                <td class="dim"><b>Our training data</b> &mdash; derelict ghost gear.</td></tr>
            <tr><td>GhostVision<br><span style="font-size:10px;color:var(--ink-3);">JMSE 14(10):951, 2025</span></td>
                <td class="dim">NOASSERTION</td>
                <td class="dim">Closest prior system. Not vendored &mdash; licence unresolved.</td></tr>
            <tr><td>PINGMapper<br><span style="font-size:10px;color:var(--ink-3);">Earth &amp; Space Science, 2022</span></td>
                <td style="color:var(--green);">MIT</td>
                <td class="dim">Pipeline shape and output conventions.</td></tr>
            <tr><td>sidescantools</td><td class="dim">GPL-3.0</td>
                <td class="dim">Studied; not vendored.</td></tr>
            <tr><td>AI4Shipwrecks<br><span style="font-size:10px;color:var(--ink-3);">arXiv 2401.14546</span></td>
                <td style="color:var(--green);">MIT</td>
                <td class="dim">Route to a wreck class later.</td></tr>
          </table>
        </div>

        <div class="box violet">
          <h3 class="sec" style="color:var(--violet);">Method references</h3>
          <ul class="b tight" style="font-size:12px;">
            <li><b>Lee (1980)</b> &mdash; adaptive speckle filtering, DOI 10.1109/TPAMI.1980.4766994</li>
            <li><b>Platt (1999)</b> &mdash; probabilistic outputs for large-margin classifiers, MIT Press</li>
            <li><b>Guo et al. (2017)</b> &mdash; on calibration of modern neural networks, arXiv:1706.04599</li>
            <li><b>Ultralytics YOLO11</b> &mdash; detection backbone (AGPL-3.0, interface-isolated)</li>
          </ul>
        </div>

        <div class="box amber t-amber">
          <h3 class="sec amber">Architectures reviewed and deliberately not adopted</h3>
          <div style="font-size:12px;line-height:1.42;color:#5A3D00;">
            <b>SSM-DETR &nbsp;&middot;&nbsp; TR-YOLOv5s &nbsp;&middot;&nbsp; MSF-DETR
            &nbsp;&middot;&nbsp; LEF-RT-DETR.</b> Each needs
            <b>2.6&times;&ndash;44&times; our compute</b> for single-digit AP gains,
            measured on forward-looking sonar or on datasets we cannot obtain to reproduce.
            Our bottleneck is false alarms on empty seabed &mdash; a bigger backbone does
            not fix that.
          </div>
        </div>
      </div>

      <div style="display:flex;flex-direction:column;gap:11px;">
        <div class="box teal">
          <h3 class="sec teal">How our approach compares</h3>
          <table class="t" style="font-size:12px;">
            <tr><th style="width:30%;">Capability</th>
                <th class="c" style="width:19%;background:var(--tint-teal);color:var(--teal);">AQUA-SHIELD</th>
                <th class="c" style="width:17%;">Manual analyst</th>
                <th class="c" style="width:16%;">Generic YOLO</th>
                <th class="c">Published SSS detectors</th></tr>
            <tr><td>Detects bottom objects in SSS</td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="yes">&#10003;</span></td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="yes">&#10003;</span></td></tr>
            <tr><td>Independent physical verification</td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="yes">&#10003;</span></td>
                <td class="c"><span class="no">&#10007;</span></td><td class="c"><span class="no">&#10007;</span></td></tr>
            <tr><td>Calibrated confidence</td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="no">&#10007;</span></td>
                <td class="c"><span class="no">&#10007;</span></td><td class="c"><span class="part">rare</span></td></tr>
            <tr><td>Refuses to invent a coordinate</td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="yes">&#10003;</span></td>
                <td class="c"><span class="no">&#10007;</span></td><td class="c"><span class="no">&#10007;</span></td></tr>
            <tr><td>Per-detection evidence &amp; audit trail</td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="part">informal</span></td>
                <td class="c"><span class="no">&#10007;</span></td><td class="c"><span class="no">&#10007;</span></td></tr>
            <tr><td>Runs offline on survey hardware</td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="yes">&#10003;</span></td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="part">varies</span></td></tr>
            <tr><td>Scales to thousands of km</td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="no">&#10007;</span></td>
                <td class="c"><span class="yes">&#10003;</span></td><td class="c"><span class="yes">&#10003;</span></td></tr>
          </table>
        </div>

        <div class="box red" style="flex:1;">
          <h3 class="sec red">What we have NOT done &mdash; stated deliberately</h3>
          <ul class="b tight" style="font-size:12px;">
            <li>Ghost-gear detection has <b>one raw training result</b> (mAP50 0.323). It
                has not yet been through the verification stage, calibration, or an
                ablation.</li>
            <li>Geolocation accuracy has <b>never been validated</b> &mdash; the licensed
                data available to us ships no navigation track.</li>
            <li>Detection <b>degrades under added speckle noise</b> (measured).</li>
            <li>Never run on a <b>Jetson or a live AUV</b>.</li>
            <li>Our own sonar preprocessing chain <b>did not help</b> when measured
                properly, and ships disabled by default.</li>
          </ul>
          <div style="font-size:12px;color:var(--ink-2);margin-top:8px;line-height:1.42;
                      border-top:1px solid var(--rule);padding-top:8px;">
            <b>Roadmap to the finals:</b> &nbsp;1 &middot; fit + ablate verification on the
            ghost-gear model &nbsp;&middot; 2 &middot; run the baseline comparison
            &nbsp;&middot; 3 &middot; report the reliability curve and ECE &nbsp;&middot;
            4 &middot; benchmark the ONNX export on Jetson &nbsp;&middot; 5 &middot;
            validate geolocation against a survey that ships a navigation track.
          </div>
        </div>
      </div>
    </div>
  </div>
</div>"""


SLIDES = [S1, S2, S3, S4, S5, S6]

HTML = f"""<!doctype html><html><head><meta charset="utf-8">
<link rel="stylesheet" href="{(DECK / 'slides.css').as_uri()}">
</head><body>{''.join(SLIDES)}</body></html>"""

page_html = BUILD / "deck.html"
page_html.write_text(HTML)
print(f"wrote {page_html}")


# --------------------------------------------------------------- render ----
CHROME_BIN = ("/Users/earther/Library/Caches/ms-playwright/chromium-1228/chrome-mac-arm64/"
              "Google Chrome for Testing.app/Contents/MacOS/Google Chrome for Testing")

from playwright.sync_api import sync_playwright  # noqa: E402

pngs = []
with sync_playwright() as pw:
    browser = (pw.chromium.launch(executable_path=CHROME_BIN)
               if Path(CHROME_BIN).exists() else pw.chromium.launch(channel="chrome"))
    page = browser.new_page(viewport={"width": 1600, "height": 900}, device_scale_factor=2)
    page.goto(page_html.as_uri(), wait_until="networkidle", timeout=90_000)
    page.wait_for_timeout(1500)

    for i, el in enumerate(page.query_selector_all(".slide"), 1):
        p = BUILD / f"slide{i}.png"
        el.screenshot(path=str(p))
        pngs.append(p)
        print(f"  rendered {p.name}")

    page.pdf(path=str(OUT_DIR / "AQUA_SHIELD_SIH_PITCH.pdf"),
             width="1600px", height="900px", print_background=True,
             margin={"top": "0", "bottom": "0", "left": "0", "right": "0"})
    browser.close()
print(f"wrote {OUT_DIR / 'AQUA_SHIELD_SIH_PITCH.pdf'}")


# ----------------------------------------------------------------- pptx ----
from pptx import Presentation                     # noqa: E402
from pptx.util import Inches                      # noqa: E402

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for p in pngs:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
prs.save(str(OUT_DIR / "AQUA_SHIELD_SIH_PITCH.pptx"))
print(f"wrote {OUT_DIR / 'AQUA_SHIELD_SIH_PITCH.pptx'}")
