#!/usr/bin/env python3
"""Fill the OFFICIAL SIH2026 PPT template with AQUA-SHIELD content.

This is the FALLBACK deck, kept in case the organisers require their exact
template file. The primary submission deck is built by
scripts/deck/build_deck.py and written to AQUA_SHIELD_SIH_PITCH.{pdf,pptx}.

Starts from the organiser-supplied template (already drafted with our project's
content by hand) and only touches text runs that are factually wrong or stale:
  - hardware benchmarks table claimed an Intel/RTX4060 Windows laptop; the
    project actually runs on Apple M5 (MPS, no CUDA) -- every number in
    docs/BENCHMARKS.md was measured on that machine, not the one named here
  - two "future work" bullets that are now stale since E14 (first-ever
    ghost-gear training run) actually completed

Everything else (official SIH branding, footer, slide numbers, layout,
fonts, the already-accurate content on slides 1/2/3/5/6) is left untouched --
the point of using the organiser's own file is exact template compliance.

The pipeline diagram on slide 3 (image3.png) is replaced with an accurate one
(scripts/build_pipeline_diagram.py) because the original graphic omitted the
detection stage entirely and showed calibration<->filtering as a loop, which
isn't how the pipeline works.
"""
import shutil
import zipfile
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
# vendored so this is reproducible -- the original arrived over WhatsApp and its
# temp file is not durable
SRC = ROOT / "assets" / "SIH2026_official_template.pptx"
OUT = ROOT / "docs" / "AQUA_SHIELD_SIH_OFFICIAL_TEMPLATE.pptx"
DIAGRAM = ROOT / "docs" / "images" / "pipeline_diagram_sih.png"

REPLACEMENTS = {
    "CPU: Intel Core i7 14650HX": "CPU: Apple M5 (10-core)",
    "GPU: NVIDIA RTX 4060 Laptop GPU": "GPU: Apple M5 GPU via MPS (no CUDA)",
    "RAM: 16 GB RAM": "RAM: 24 GB unified memory",
    "Future validation using ghost-gear datasets":
        "Ghost-gear model trained (E14): mAP50 0.32, first result, pending full verification",
    "Limited labelled ghost-gear data":
        "Ghost-gear result is one raw run, not yet ablated or verified",
}

prs = Presentation(str(SRC))

applied = {k: 0 for k in REPLACEMENTS}
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in REPLACEMENTS.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        applied[old] += 1

for old, count in applied.items():
    status = "OK" if count == 1 else f"WARNING: matched {count} times"
    print(f"[{status}] '{old}'")

# emphasise the winning row of the ablation table (slide 4) -- best false-alarm
# reduction while keeping the most true positives of any filtered configuration
s4 = prs.slides[3]
for shape in s4.shapes:
    if shape.has_table and shape.table.cell(3, 0).text == "Detector + Learned FP Filter":
        for c in range(3):
            for para in shape.table.cell(3, c).text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True
        print("[OK] bolded 'Detector + Learned FP Filter' row in ablation table")

prs.save(str(OUT))

# swap the inaccurate placeholder pipeline graphic for the corrected one, in place
with zipfile.ZipFile(OUT, "r") as zin:
    names = zin.namelist()
    data = {n: zin.read(n) for n in names}

target = next(n for n in names if n.endswith("media/image3.png"))
data[target] = DIAGRAM.read_bytes()

tmp = OUT.with_suffix(".tmp.pptx")
with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for n in names:
        zout.writestr(n, data[n])
shutil.move(str(tmp), str(OUT))
print(f"replaced {target} with corrected pipeline diagram")
print(f"wrote {OUT}")
