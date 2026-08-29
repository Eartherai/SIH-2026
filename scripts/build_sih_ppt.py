#!/usr/bin/env python3
"""Build the SIH 2026 submission deck (PS 26057) from measured project facts.

Every number here is pulled from docs/BENCHMARKS.md and experiments/registry.jsonl
at the time this was written. If a number changes, regenerate this deck --
do not hand-edit numbers into the .pptx.

Follows the official SIH 6-slide template:
  1. Problem & Proposed Solution (title info in the banner)
  2. Idea & Innovation (USP)
  3. Technical Approach (architecture)
  4. Feasibility & Viability
  5. Impact & Benefits
  6. Research & References
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "AQUA_SHIELD_SIH_PITCH.pptx"
IMG_VERIFY = ROOT / "docs" / "images" / "verification_effect.png"
IMG_FAIL = ROOT / "docs" / "images" / "failure_gallery.png"

# ---------------------------------------------------------------- palette --
NAVY_DARK = RGBColor(0x08, 0x14, 0x22)
NAVY = RGBColor(0x0B, 0x1E, 0x33)
NAVY_MID = RGBColor(0x11, 0x2A, 0x44)
NAVY_LIGHT = RGBColor(0x16, 0x3A, 0x5C)
CYAN = RGBColor(0x2D, 0xE0, 0xD8)
CYAN_DIM = RGBColor(0x5B, 0x9E, 0xA8)
AMBER = RGBColor(0xF5, 0xB3, 0x41)
WHITE = RGBColor(0xF2, 0xF7, 0xFA)
GRAY = RGBColor(0x9F, 0xB3, 0xC8)
GRAY_DIM = RGBColor(0x6E, 0x84, 0x99)
GREEN = RGBColor(0x5D, 0xD9, 0x8B)
RED = RGBColor(0xE0, 0x6C, 0x6C)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers --
def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    return s


def no_shadow(shape):
    shape.shadow.inherit = False


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    no_shadow(shape)


def add_rect(slide, x, y, w, h, color, line_color=None, line_w=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    if line_color:
        r.line.color.rgb = line_color
        r.line.width = line_w or Pt(0.75)
    else:
        r.line.fill.background()
    no_shadow(r)
    return r


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, line_spacing=None):
    """runs: list of paragraphs, each a list of (text, size, color, bold, italic) tuples,
    OR a single tuple list for a single paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for text, size, color, bold, italic in para_runs:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = FONT
    return tb


def bullets(slide, x, y, w, h, items, size=13.5, color=WHITE, gap=6, bullet_color=CYAN,
            leading=1.12):
    """items: list of (text, level, style) where style in {'normal','honest','good'}."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        text, level, style = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = leading
        marker = "▸  " if level == 0 else "–  "
        r0 = p.add_run()
        r0.text = marker
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color if level == 0 else GRAY_DIM
        r0.font.name = FONT
        r0.font.bold = level == 0
        if level == 1:
            p.level = 0
            r0.text = "    " + r0.text
        r1 = p.add_run()
        r1.text = text
        r1.font.size = Pt(size - (0 if level == 0 else 0.5))
        r1.font.name = FONT
        if style == "honest":
            r1.font.color.rgb = AMBER
            r1.font.italic = False
        elif style == "good":
            r1.font.color.rgb = GREEN
        else:
            r1.font.color.rgb = color if level == 0 else GRAY
        r1.font.bold = style == "good"
    return tb


def header(slide, kicker, title):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), NAVY_DARK)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Pt(2.2), CYAN)
    add_text(slide, Inches(0.55), Inches(0.14), Inches(9.0), Inches(0.32),
              [[(kicker, 12, CYAN, True, False)]])
    add_text(slide, Inches(0.55), Inches(0.44), Inches(11.5), Inches(0.62),
              [[(title, 26, WHITE, True, False)]])


def footer(slide, n):
    add_text(slide, Inches(0.55), Inches(7.18), Inches(6), Inches(0.28),
              [[("AQUA-SHIELD", 9.5, CYAN_DIM, True, False),
                ("  ·  PS 26057  ·  MoES / NIOT  ·  Disaster Management", 9.5, GRAY_DIM, False, False)]])
    add_text(slide, Inches(12.2), Inches(7.18), Inches(0.6), Inches(0.28),
              [[(f"{n}/6", 9.5, GRAY_DIM, False, False)]], align=PP_ALIGN.RIGHT)


def stat_card(slide, x, y, w, h, value, label, value_color=CYAN, value_size=22):
    card = add_rect(slide, x, y, w, h, NAVY_MID, line_color=NAVY_LIGHT, line_w=Pt(1))
    add_text(slide, x + Inches(0.12), y + Inches(0.08), w - Inches(0.24), h - Inches(0.55),
              [[(value, value_size, value_color, True, False)]], anchor=MSO_ANCHOR.BOTTOM)
    add_text(slide, x + Inches(0.12), y + h - Inches(0.42), w - Inches(0.24), Inches(0.38),
              [[(label, 10.5, GRAY, False, False)]], anchor=MSO_ANCHOR.TOP)


def styled_table(slide, x, y, w, h, header_row, rows, col_widths, font_size=11.5,
                  header_size=11.5, row_h=None, highlight_last_col=False):
    n_rows = len(rows) + 1
    n_cols = len(header_row)
    gtable = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = gtable.table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw
    if row_h:
        for r in range(n_rows):
            tbl.rows[r].height = row_h

    def style_cell(cell, text, bg, fg, bold=False, size=font_size, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.margin_left = Pt(6)
        cell.margin_right = Pt(6)
        cell.margin_top = Pt(3)
        cell.margin_bottom = Pt(3)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(text)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = fg
        r.font.name = FONT

    for c, htext in enumerate(header_row):
        style_cell(tbl.cell(0, c), htext, NAVY_LIGHT, CYAN, bold=True, size=header_size)
    for ri, row in enumerate(rows):
        bg = NAVY_MID if ri % 2 == 0 else NAVY
        for c, val in enumerate(row):
            fg = WHITE
            bold = False
            if highlight_last_col and c == len(row) - 1:
                fg = CYAN
                bold = True
            style_cell(tbl.cell(ri + 1, c), val, bg, fg, bold=bold)
    # kill default pptx table style banding
    tbl.first_row = False
    tbl.horz_banding = False
    return gtable


def pill(slide, x, y, text, color=CYAN, text_color=NAVY_DARK, w=None, h=Inches(0.34), size=11):
    w = w or Inches(0.4 + 0.095 * len(text))
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.adjustments[0] = 0.5
    set_fill(p, color)
    tf = p.text_frame
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = 0; tf.margin_bottom = 0
    pr = tf.paragraphs[0]
    pr.alignment = PP_ALIGN.CENTER
    r = pr.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    r.font.name = FONT
    return p, w


# ============================================================== SLIDE 1 ===
s = add_slide()
add_rect(s, 0, 0, SLIDE_W, Inches(1.9), NAVY_DARK)
add_rect(s, 0, Inches(1.9), SLIDE_W, Pt(2.2), CYAN)
add_text(s, Inches(0.55), Inches(0.18), Inches(9), Inches(0.3),
          [[("SMART INDIA HACKATHON 2026", 12, CYAN, True, False)]])
add_text(s, Inches(0.55), Inches(0.48), Inches(11.5), Inches(0.85),
          [[("AQUA-SHIELD", 40, WHITE, True, False)]])
add_text(s, Inches(0.58), Inches(1.32), Inches(11.5), Inches(0.5),
          [[("Detection → Verification → Localization → Action", 15, CYAN, False, True)]])
tags_y = Inches(2.15)
x = Inches(0.55)
for label in ["PS 26057", "MoES / NIOT", "Software", "Disaster Management"]:
    _, w = pill(s, x, tags_y, label, color=NAVY_LIGHT, text_color=WHITE, size=11.5)
    x += w + Inches(0.16)

add_text(s, Inches(0.55), Inches(2.75), Inches(7.0), Inches(0.4),
          [[("PROBLEM", 13, AMBER, True, False)]])
bullets(s, Inches(0.55), Inches(3.15), Inches(6.9), Inches(3.6), [
    ("Ghost nets and derelict fishing gear keep killing after they are lost — entangling marine life, smothering reefs, fouling propellers.", 0, "normal"),
    ("Finding them today means a human reading thousands of kilometres of side-scan sonar (SSS) imagery by eye.", 0, "normal"),
    ("The real difficulty is not detection.", 0, "normal"),
    ("74% of sonar frames contain no target at all — a system that fires on a fraction of those is worse than useless. The analyst stops trusting it.", 1, "honest"),
    ("Precision, not recall, is the binding constraint. That governs every design choice in this project.", 0, "good"),
], size=13.5, gap=10)

add_text(s, Inches(7.75), Inches(2.75), Inches(5.0), Inches(0.4),
          [[("OUR SOLUTION", 13, CYAN, True, False)]])
bullets(s, Inches(7.75), Inches(3.15), Inches(4.7), Inches(3.5), [
    ("A local-first pipeline: Detect → Verify against physical evidence → Calibrate confidence → Deduplicate repeat sightings → Geolocate (or refuse) → Prioritise → Export.", 0, "normal"),
    ("Runs fully offline on a single laptop — no cloud, no per-image cost, no data leaving the vessel.", 0, "normal"),
    ("Built for two datasets so far: naval mine-like objects (MILCO/NOMBO) and derelict crab-pot ghost gear.", 0, "normal"),
], size=13, gap=10)
footer(s, 1)

# ============================================================== SLIDE 2 ===
s = add_slide()
header(s, "IDEA & INNOVATION", "What's actually ours")
add_text(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.55),
          [[("What we do NOT claim: ", 13, AMBER, True, False),
            ("sonar debris detection itself is not novel — GhostVision (JMSE 2025) already does detection + georeferencing for derelict gear. We say so in our own repository.", 13, GRAY, False, True)]])

cards = [
    ("01", "A learned false-positive filter over physical evidence",
     "Ten features measured directly from pixels — shadow coherence, contrast, compactness, relative texture — independent of the detector's own opinion. Fitted on a held-out survey, with per-detection attribution."),
    ("02", "Confidence that means something",
     "Platt calibration fitted on held-out data. When it isn't fitted, every hazard is stamped calibrated: false — never a silently-wrong number."),
    ("03", "Refusal as a feature",
     "No navigation metadata → no coordinate. null, plus the reason. A fabricated latitude exports cleanly to CSV and sends a cleanup vessel to open water."),
    ("04", "Confidence ≠ Priority",
     "“Is it real?” and “should you care?” are treated as two different questions, scored separately."),
]
cx, cy, cw, ch, gap = Inches(0.55), Inches(2.05), Inches(3.62), Inches(2.15), Inches(0.18)
for i, (num, title, body) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = cx + col * (cw + gap)
    y = cy + row * (ch + gap)
    card = add_rect(s, x, y, cw, ch, NAVY_MID, line_color=NAVY_LIGHT, line_w=Pt(1))
    add_text(s, x + Inches(0.18), y + Inches(0.12), Inches(0.8), Inches(0.4),
              [[(num, 20, CYAN, True, False)]])
    add_text(s, x + Inches(0.18), y + Inches(0.52), cw - Inches(0.36), Inches(0.55),
              [[(title, 13.5, WHITE, True, False)]])
    add_text(s, x + Inches(0.18), y + Inches(1.08), cw - Inches(0.36), Inches(1.0),
              [[(body, 10.8, GRAY, False, False)]])

# right column: the diagnostic-bug story
rx = Inches(8.35)
box = add_rect(s, rx, Inches(2.05), Inches(4.42), Inches(4.5), NAVY_DARK, line_color=AMBER, line_w=Pt(1.25))
add_text(s, rx + Inches(0.2), Inches(2.2), Inches(4.0), Inches(0.35),
          [[("ITS WEIGHTS CAUGHT A BUG IN OUR OWN PIPELINE", 11.5, AMBER, True, False)]])
bullets(s, rx + Inches(0.2), Inches(2.65), Inches(4.0), Inches(3.7), [
    ("An early fit gave acoustic shadow a large NEGATIVE weight — the opposite of the physics.", 0, "normal"),
    ("Root cause: a preprocessing chain applied at inference that the detector had never been trained on.", 0, "normal"),
    ("Fixing the mismatch recovered a 12× F1 improvement (0.012 → 0.144) and turned both shadow features positive again.", 0, "good"),
], size=11.8, gap=12)
add_text(s, rx + Inches(0.2), Inches(5.55), Inches(4.0), Inches(0.9),
          [[("An inspectable verification stage is a diagnostic instrument, not just a classifier — a hand-tuned threshold would have hidden the same defect silently.", 11, CYAN, False, True)]])
footer(s, 2)

# ============================================================== SLIDE 3 ===
s = add_slide()
header(s, "TECHNICAL APPROACH", "Architecture & stack")

# pipeline diagram
stages = ["RAW\nSONAR", "QC", "PRE-\nPROCESS\n(off)", "TILING", "DETECTION\nYOLO11n",
          "FP\nFILTER", "CALI-\nBRATION", "DEDUP", "GEO-\nLOCATE", "PRIORITY", "REPORT"]
n = len(stages)
diag_y = Inches(1.5)
box_w = Inches(1.02)
box_h = Inches(0.72)
gap_w = Inches(0.075)
total_w = n * box_w + (n - 1) * gap_w
start_x = (SLIDE_W - total_w) / 2
xs = []
for i, label in enumerate(stages):
    x = start_x + i * (box_w + gap_w)
    xs.append(x)
    fill = NAVY_LIGHT if label != "PRE-\nPROCESS\n(off)" else NAVY_DARK
    line = CYAN if label != "PRE-\nPROCESS\n(off)" else AMBER
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, diag_y, box_w, box_h)
    b.adjustments[0] = 0.12
    b.fill.solid(); b.fill.fore_color.rgb = fill
    b.line.color.rgb = line; b.line.width = Pt(1.25)
    no_shadow(b)
    tf = b.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(8.3)
    r.font.bold = True
    r.font.color.rgb = AMBER if line == AMBER else WHITE
    r.font.name = FONT
    if i < n - 1:
        arrow_x = x + box_w
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, arrow_x, diag_y + box_h / 2,
                                       arrow_x + gap_w, diag_y + box_h / 2)
        conn.line.color.rgb = CYAN_DIM
        conn.line.width = Pt(1.25)

add_text(s, start_x, diag_y + box_h + Inches(0.08), total_w, Inches(0.3),
          [[("preprocessing implemented, measured, disabled by default — it made detection worse", 9.5, AMBER, False, True)]],
          align=PP_ALIGN.CENTER)

rows = [
    ["QC", "Dynamic range, speckle index, dropout rows, water-column detection"],
    ["Preprocessing", "Dropout repair · water-column removal · Lee speckle filter · gain normalisation — off by default (measured, see Slide 5)"],
    ["Tiling", "Overlapping tiles at native resolution; seam duplicates merged by IoU or intersection-over-smaller"],
    ["Detection", "YOLO11n — 2.58M params, 6.3 GFLOPs, backend-swappable"],
    ["Verification", "10 physical features → logistic model fitted on held-out survey"],
    ["Geolocation", "GeoTIFF affine · per-ping navigation · or refuse with a stated reason"],
]
styled_table(s, Inches(0.55), Inches(2.95), Inches(7.15), Inches(3.35),
             ["Stage", "Substance"], rows,
             [Inches(1.55), Inches(5.6)], font_size=10.8, header_size=11)

add_text(s, Inches(8.05), Inches(2.95), Inches(4.7), Inches(0.35),
          [[("TECHNOLOGY STACK", 12, CYAN, True, False)]])
stack_items = ["Python 3.12", "PyTorch 2.13 (MPS · CPU fallback, no CUDA assumed)",
               "OpenCV", "pyproj", "Streamlit", "FastAPI", "SQLite", "ONNX Runtime"]
y = Inches(3.35)
for item in stack_items:
    pill(s, Inches(8.05), y, item, color=NAVY_LIGHT, text_color=WHITE, size=10.5, h=Inches(0.34))
    y += Inches(0.42)
footer(s, 3)

# ============================================================== SLIDE 4 ===
s = add_slide()
header(s, "FEASIBILITY & VIABILITY", "It already runs — measured, not projected")

stat_card(s, Inches(0.55), Inches(1.4), Inches(1.95), Inches(1.05), "21.4 ms", "MPS inference / tile")
stat_card(s, Inches(2.62), Inches(1.4), Inches(1.95), Inches(1.05), "3.8×", "MPS speed-up vs CPU")
stat_card(s, Inches(4.69), Inches(1.4), Inches(1.95), Inches(1.05), "37.4 /s", "Throughput, frames")
stat_card(s, Inches(6.76), Inches(1.4), Inches(1.95), Inches(1.05), "640 MB", "Peak memory")
stat_card(s, Inches(8.83), Inches(1.4), Inches(1.95), Inches(1.05), "10.6 MB", "ONNX export size")
stat_card(s, Inches(10.90), Inches(1.4), Inches(1.88), Inches(1.05), "8.5 ms", "ONNX inference")
add_text(s, Inches(0.55), Inches(2.55), Inches(12.2), Inches(0.3),
          [[("Measured on an Apple M5, 24GB unified memory — a laptop, not a cluster.", 10.5, GRAY_DIM, False, True)]])

add_text(s, Inches(0.55), Inches(3.0), Inches(6.0), Inches(0.35),
          [[("ABLATION — 612 held-out frames, 473 target-free", 12, CYAN, True, False)]])
rows = [
    ["Detector only", "0.247", "21", "37 / 473"],
    ["+ hand-written rules", "0.300", "12", "18 / 473"],
    ["+ learned FP filter", "0.322", "19", "25 / 473"],
]
styled_table(s, Inches(0.55), Inches(3.4), Inches(6.5), Inches(1.55),
             ["Configuration", "Precision", "True positives", "Falsely-alarmed frames"], rows,
             [Inches(2.1), Inches(1.4), Inches(1.5), Inches(1.5)], font_size=10.8, header_size=10.8,
             highlight_last_col=False)

add_text(s, Inches(0.55), Inches(5.15), Inches(6.5), Inches(1.9),
          [[("Data legitimacy — ", 11.5, CYAN, True, False),
            ("MILCO/NOMBO real AUV side-scan, CC BY 4.0, DOI 10.6084/m9.figshare.24574879. Derelict crab-pot ghost gear, CC BY-SA 4.0 (PINGEcosystem). Every dependency licence verified programmatically.", 11.5, GRAY, False, False)],
           [("Evaluated honestly — ", 11.5, CYAN, True, False),
            ("split by acquisition year / recording ID, never randomly. Train 2015+2010 → calibrate 2017 → test 2018+2021. A random split leaks: consecutive frames share seabed, gain settings, and often the same object.", 11.5, GRAY, False, False)]],
          line_spacing=1.15)

rx = Inches(7.35)
box = add_rect(s, rx, Inches(3.0), Inches(5.45), Inches(4.05), NAVY_DARK, line_color=NAVY_LIGHT, line_w=Pt(1))
add_text(s, rx + Inches(0.2), Inches(3.15), Inches(5.0), Inches(0.35),
          [[("RISKS & MITIGATION", 12, CYAN, True, False)]])
bullets(s, rx + Inches(0.2), Inches(3.55), Inches(5.05), Inches(3.4), [
    ("Coastal SSS training data is scarce.", 0, "normal"),
    ("Mitigated with leakage-free, recording-level splitting and honest reporting rather than inflated splits.", 1, "normal"),
    ("Detection backend (Ultralytics) is AGPL-3.0.", 0, "normal"),
    ("Architecture keeps the detector backend swappable; ONNX export path is backend-independent.", 1, "normal"),
    ("Compute is constrained on AUV/edge hardware.", 0, "normal"),
    ("ONNX export measured at 10.6MB / 8.5ms — sized for edge deployment, not just the dev laptop.", 1, "normal"),
    ("Detection collapses under added speckle noise (measured).", 0, "honest"),
    ("Alternative noise-robust checkpoint trained and shipped separately, honestly labelled as a trade-off, not silently swapped in.", 1, "normal"),
], size=11.2, gap=8)
footer(s, 4)

# ============================================================== SLIDE 5 ===
s = add_slide()
header(s, "IMPACT & BENEFITS", "Measured outcomes, not projections")

stat_card(s, Inches(0.55), Inches(1.4), Inches(2.85), Inches(1.15), "37 → 25", "Falsely-alarmed frames / 473 (−32%)", value_color=GREEN)
stat_card(s, Inches(3.55), Inches(1.4), Inches(2.85), Inches(1.15), "0.247 → 0.322", "Verification-stage precision (+30%)", value_color=GREEN)
stat_card(s, Inches(6.55), Inches(1.4), Inches(2.85), Inches(1.15), "19 / 21", "True positives kept while cutting alarms", value_color=CYAN)
stat_card(s, Inches(9.55), Inches(1.4), Inches(3.23), Inches(1.15), "0.323 mAP50", "First-ever ghost-gear result (E14, new)", value_color=AMBER, value_size=19)

add_text(s, Inches(9.55), Inches(2.65), Inches(3.23), Inches(0.65),
          [[("Raw training number — not yet through the verification/calibration stage. Reported as-is, not oversold.", 9, AMBER, False, True)]])

add_text(s, Inches(0.55), Inches(2.75), Inches(8.6), Inches(0.35),
          [[("OPERATIONAL BENEFIT", 12.5, CYAN, True, False)]])
bullets(s, Inches(0.55), Inches(3.15), Inches(8.6), Inches(1.9), [
    ("The analyst reads a ranked hazard register instead of scrolling raw sonar imagery.", 0, "normal"),
    ("Output lands in QGIS as GeoJSON — a cleanup vessel can be tasked directly from the report.", 0, "normal"),
    ("Runs fully offline on survey hardware: no cloud, no per-image cost, no data leaving the ship.", 0, "normal"),
    ("Deduplication merges repeat sightings into unique hazards; positional uncertainty tightens by ~√N over repeat fixes.", 0, "normal"),
], size=12.5, gap=9)

add_text(s, Inches(0.55), Inches(5.15), Inches(8.6), Inches(0.35),
          [[("WHO THIS SERVES", 12.5, CYAN, True, False)]])
bullets(s, Inches(0.55), Inches(5.55), Inches(8.6), Inches(1.6), [
    ("MoES / NIOT survey and cleanup operations — a direct fit for the problem statement's own agency.", 0, "normal"),
    ("Fisheries and coastal management — derelict gear (\"ghost fishing\") keeps killing catch and habitat long after it's lost.", 0, "normal"),
    ("We do not claim a percentage of analyst time saved — that needs a user study we have not run. What's shown here is a measured drop in false alarms and a measured processing rate.", 0, "honest"),
], size=12, gap=9)
footer(s, 5)

# ============================================================== SLIDE 6 ===
s = add_slide()
header(s, "RESEARCH & REFERENCES", "Prior art, and what we take from it")

rows = [
    ["GhostVision (JMSE 14(10):951, 2025)", "NOASSERTION", "Closest system. Not vendored — licence unresolved."],
    ["PINGMapper (Earth & Space Science, 2022)", "MIT", "Pipeline shape and output conventions"],
    ["sidescantools", "GPL-3.0", "Studied; not vendored"],
    ["AI4Shipwrecks (arXiv 2401.14546)", "MIT", "Route to a wreck class later"],
    ["MILCO/NOMBO (Data in Brief 53:110132)", "CC BY 4.0", "Training data — naval mine-like objects"],
    ["sss-crab-pot-detection-ds (PINGEcosystem)", "CC BY-SA 4.0", "Training data — derelict ghost gear (E14, new)"],
]
styled_table(s, Inches(0.55), Inches(1.45), Inches(8.05), Inches(2.85),
             ["Work", "Licence", "Relationship"], rows,
             [Inches(3.35), Inches(1.35), Inches(3.35)], font_size=10.3, header_size=10.8, row_h=Inches(0.44))

add_text(s, Inches(0.55), Inches(4.5), Inches(8.05), Inches(0.32),
          [[("METHOD REFERENCES", 11.5, CYAN, True, False)]])
add_text(s, Inches(0.55), Inches(4.85), Inches(8.05), Inches(0.6),
          [[("Lee (1980) adaptive speckle filtering  ·  Platt (1999) probabilistic outputs  ·  Guo et al. (2017) calibration of modern neural networks", 10.8, GRAY, False, False)]])

box = add_rect(s, Inches(0.55), Inches(5.55), Inches(8.05), Inches(1.5), NAVY_DARK, line_color=AMBER, line_w=Pt(1.25))
add_text(s, Inches(0.75), Inches(5.68), Inches(7.6), Inches(0.3),
          [[("WHAT WE HAVE NOT DONE — STATED DELIBERATELY", 11, AMBER, True, False)]])
bullets(s, Inches(0.75), Inches(6.02), Inches(7.6), Inches(1.0), [
    ("Ghost-gear detection has one raw training result (mAP50 0.32) — not yet run through verification, calibration, or full ablation.", 0, "honest"),
    ("Geolocation accuracy never validated (our data ships no navigation); detection degrades under added speckle noise; never run on a Jetson or an AUV; detector backend is AGPL-3.0.", 0, "honest"),
], size=10.3, gap=5)

# right: failure gallery image, honest visual
from PIL import Image
img = Image.open(IMG_FAIL)
iw, ih = img.size
disp_w = Inches(4.15)
disp_h = Emu(int(disp_w * ih / iw))
ix = Inches(13.03) - disp_w
s.shapes.add_picture(str(IMG_FAIL), ix, Inches(1.45), width=disp_w)
add_text(s, ix, Inches(1.45) + disp_h + Inches(0.06), disp_w, Inches(0.35),
          [[("measured failure modes, shown not hidden", 9, GRAY_DIM, False, True)]], align=PP_ALIGN.CENTER)
footer(s, 6)

prs.save(str(OUT))
print(f"wrote {OUT}")
